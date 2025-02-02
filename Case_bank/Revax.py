import os
import json
import time
from pptx import Presentation
from mistralai import Mistral

# Mistral API Setup
api_key = os.environ.get("MISTRAL_API_KEY", "9zMRdnt9VJz3fprSFW4ydGkKmC2sdYMF")
model = "mistral-embed"
client = Mistral(api_key=api_key)

# Paths
base_dir = "/var/www/html/Case_bank"
target_dir = os.path.join(base_dir, "Target")
processed_dir = os.path.join(base_dir, "Processed")

def extract_text_from_pptx(file_path):
    """Extract text from a .pptx file and return it as a list of chunks with metadata."""
    text_chunks = []
    presentation = Presentation(file_path)
    for slide_number, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip() != "":
                text_chunks.append({
                    "slide_number": slide_number + 1,
                    "text": shape.text.strip()
                })
    return text_chunks

def embed_text_chunks(text_chunks):
    """Embed text chunks using Mistral API with rate limiting."""
    embedded_text = []
    for chunk in text_chunks:
        try:
            response = client.embeddings.create(model=model, inputs=[chunk["text"]])
            embedded_text.append({
                "slide_number": chunk["slide_number"],
                "text": chunk["text"],
                "embedding": response.data[0].embedding
            })
        except Exception as e:
            print(f"Error embedding text: {e}")
            time.sleep(2)  # Wait 1 second before retrying
    return embedded_text

def process_pptx(file_path):
    """Process a single .pptx file, extracting and embedding its text."""
    print(f"Processing: {file_path}")
    text_chunks = extract_text_from_pptx(file_path)
    if not text_chunks:
        print(f"No text found in: {file_path}")
        return
    
    embedded_text = embed_text_chunks(text_chunks)
    output_file = os.path.join(processed_dir, os.path.basename(file_path) + ".json")
    with open(output_file, "w") as f:
        json.dump({
            "file_name": os.path.basename(file_path),
            "embedded_text": embedded_text
        }, f, indent=2)
    print(f"Processed and saved: {output_file}")

def main():
    """Process all .pptx files in the target directory."""
    if not os.path.exists(processed_dir):
        os.makedirs(processed_dir)
    
    for file_name in os.listdir(target_dir):
        if file_name.endswith(".pptx"):
            file_path = os.path.join(target_dir, file_name)
            process_pptx(file_path)

if __name__ == "__main__":
    main()
